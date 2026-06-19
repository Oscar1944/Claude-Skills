#!/usr/bin/env python3
"""
FastAPI Web Service for Document Extraction
"""

import os
import tempfile
import pdfplumber
from pathlib import Path
from fastapi import FastAPI, File, UploadFile, HTTPException
import uvicorn

app = FastAPI(title="Document Extraction API")


@app.get("/")
async def root():
    return {"status": "ok", "endpoints": ["/extract-pdf", "/extract-pptx"]}


@app.post("/extract-pdf")
async def extract_pdf(file: UploadFile = File(...)):
    if not file.filename.lower().endswith(".pdf"):
        raise HTTPException(status_code=400, detail="File must be a PDF")

    with tempfile.NamedTemporaryFile(delete=False, suffix=".pdf") as tmp:
        tmp.write(await file.read())
        tmp_path = tmp.name

    try:
        text_pages = []
        with pdfplumber.open(tmp_path) as pdf:
            num_pages = len(pdf.pages)
            for page in pdf.pages:
                text = page.extract_text() or ""
                text_pages.append(text)

        full_text = "\n\n".join(text_pages)

        return {
            "status": "success",
            "filename": file.filename,
            "pages": num_pages,
            "text_length": len(full_text),
            "content": full_text,
        }
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))
    finally:
        os.unlink(tmp_path)


@app.post("/extract-pptx")
async def extract_pptx(file: UploadFile = File(...)):
    if not file.filename.lower().endswith(".pptx"):
        raise HTTPException(status_code=400, detail="File must be a PPTX")

    with tempfile.NamedTemporaryFile(delete=False, suffix=".pptx") as tmp:
        tmp.write(await file.read())
        tmp_path = tmp.name

    try:
        from pptx import Presentation

        prs = Presentation(tmp_path)
        slides_content = []

        for i, slide in enumerate(prs.slides, start=1):
            texts = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        t = para.text.strip()
                        if t:
                            texts.append(t)
            slides_content.append({"slide_number": i, "text": "\n".join(texts)})

        all_text = "\n\n".join(
            f"Slide {s['slide_number']}: {s['text']}" for s in slides_content
        )

        return {
            "status": "success",
            "filename": file.filename,
            "slides": len(slides_content),
            "content_length": len(all_text),
            "content": all_text,
            "slides_content": slides_content,
        }
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))
    finally:
        os.unlink(tmp_path)


if __name__ == "__main__":
    port = int(os.getenv("PORT", 8000))
    uvicorn.run(app, host="0.0.0.0", port=port)
