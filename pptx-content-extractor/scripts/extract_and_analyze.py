#!/usr/bin/env python3
"""
PPTX Content Extractor Skill
Extracts content from PPTX presentations, answers questions about content, and generates new PPTX files.
Uses FastAPI backend for file processing.
"""

import os
import sys
import json
import requests
from pathlib import Path
from typing import Dict, List, Optional, Tuple
import logging

# Configure logging
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

# Import dependencies (will be installed via requirements.txt)
try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN
    from PIL import Image
except ImportError as e:
    logger.error(f"Missing required dependency: {e}")
    sys.exit(1)


class APIClient:
    """Client for calling the FastAPI backend."""

    def __init__(self, api_url: str = None):
        if api_url is None:
            api_url = os.getenv("PPTX_EXTRACTION_API_URL", "https://claude-skill.zeabur.app")
        self.api_url = api_url
        self.extract_endpoint = f"{api_url}/extract-pptx"

    def extract_pptx(self, file_path: str) -> Dict:
        """Extract PPTX content via API."""
        try:
            pptx_file = Path(file_path)
            if not pptx_file.exists():
                return {"status": "error", "message": f"File not found: {file_path}"}

            if not pptx_file.suffix.lower() == ".pptx":
                return {"status": "error", "message": "File must be a PPTX"}

            with open(pptx_file, "rb") as f:
                files = {"file": (pptx_file.name, f, "application/vnd.openxmlformats-officedocument.presentationml.presentation")}
                response = requests.post(self.extract_endpoint, files=files, timeout=300)

            if response.status_code == 200:
                return response.json()
            else:
                return {
                    "status": "error",
                    "message": f"API error: {response.status_code}",
                    "detail": response.text
                }
        except requests.exceptions.ConnectionError:
            return {
                "status": "error",
                "message": f"Cannot connect to API at {self.api_url}",
                "hint": "Make sure the FastAPI server is running"
            }
        except Exception as e:
            return {"status": "error", "message": str(e)}

    def health_check(self) -> bool:
        """Check if API is available."""
        try:
            response = requests.get(f"{self.api_url}/", timeout=5)
            return response.status_code == 200
        except:
            return False


class PPTXReader:
    """Handles PPTX reading and content extraction via API."""

    def __init__(self, file_path: str, api_url: str = None):
        """
        Initialize PPTX reader.

        Args:
            file_path: Path to PPTX file
            api_url: URL of the FastAPI backend
        """
        if api_url is None:
            api_url = os.getenv("PPTX_EXTRACTION_API_URL", "https://claude-skill.zeabur.app")
        self.file_path = Path(file_path)
        self.api_client = APIClient(api_url)
        self.slides_content = None
        self.metadata = None
        self.api_response = None

    def read_pptx(self) -> Dict:
        """Read PPTX and extract basic information via API."""
        try:
            self.api_response = self.api_client.extract_pptx(str(self.file_path))

            if self.api_response.get("status") == "success":
                self.slides_content = self.api_response.get("slides_content", [])
                self.metadata = {
                    "filename": self.api_response.get("filename", ""),
                    "slides": self.api_response.get("slides", 0),
                }
                return {
                    "status": "success",
                    "filename": self.api_response.get("filename", ""),
                    "slides": self.api_response.get("slides", 0),
                    "content_length": self.api_response.get("content_length", 0),
                }
            else:
                error_msg = self.api_response.get("message", "Unknown error")
                logger.error(f"Error reading PPTX: {error_msg}")
                return {"status": "error", "message": error_msg}
        except Exception as e:
            logger.error(f"Error reading PPTX: {e}")
            return {"status": "error", "message": str(e)}

    def extract_slides(self) -> List[Dict]:
        """Get extracted slides content."""
        if self.slides_content is None:
            self.read_pptx()
        return self.slides_content or []

    def extract_notes(self) -> List[str]:
        """Extract speaker notes from slides."""
        if self.slides_content is None:
            self.read_pptx()

        notes = []
        for slide in self.slides_content:
            if slide.get("notes"):
                notes.append(f"Slide {slide['slide_number']}: {slide['notes']}")

        return notes

    def get_slide_count(self) -> int:
        """Get total number of slides."""
        if self.metadata is None:
            self.read_pptx()
        return self.metadata.get("slides", 0) if self.metadata else 0


class DocumentAnalyzer:
    """Analyzes presentation content and answers questions."""

    def __init__(self, slides_content: List[Dict]):
        """
        Initialize analyzer with slides content.

        Args:
            slides_content: List of extracted slides data
        """
        self.slides_content = slides_content
        self.api_key = os.getenv("ANTHROPIC_API_KEY")

    def answer_question(self, question: str) -> str:
        """Answer question based on presentation content."""
        try:
            from anthropic import Anthropic

            client = Anthropic(api_key=self.api_key)

            presentation_text = "\n\n".join(
                [f"Slide {s['slide_number']}:\n{s['text']}" for s in self.slides_content]
            )

            system_prompt = """You are a helpful assistant that answers questions based on provided presentation content.
Answer based only on the information in the presentation. If the information is not in the presentation, say so clearly.
Provide detailed, well-structured answers."""

            message = client.messages.create(
                model="claude-opus-4-8",
                max_tokens=1024,
                system=system_prompt,
                messages=[
                    {
                        "role": "user",
                        "content": f"Presentation content:\n{presentation_text}\n\nQuestion: {question}"
                    }
                ]
            )

            return message.content[0].text
        except Exception as e:
            logger.error(f"Error answering question: {e}")
            return f"Error: {str(e)}"

    def summarize(self, max_words: Optional[int] = None) -> str:
        """Generate summary of presentation."""
        try:
            from anthropic import Anthropic

            client = Anthropic(api_key=self.api_key)

            presentation_text = "\n\n".join(
                [f"Slide {s['slide_number']}:\n{s['text']}" for s in self.slides_content]
            )

            word_constraint = f"in approximately {max_words} words" if max_words else ""
            system_prompt = f"""You are a skilled summarizer. Create a comprehensive summary of the provided presentation {word_constraint}.
Include main topics, key arguments, and important conclusions. Structure the summary clearly."""

            message = client.messages.create(
                model="claude-opus-4-8",
                max_tokens=2048,
                system=system_prompt,
                messages=[
                    {
                        "role": "user",
                        "content": f"Please summarize this presentation:\n\n{presentation_text}"
                    }
                ]
            )

            return message.content[0].text
        except Exception as e:
            logger.error(f"Error summarizing: {e}")
            return f"Error: {str(e)}"

    def extract_key_points(self, num_points: int = 5) -> List[str]:
        """Extract key points from presentation."""
        try:
            from anthropic import Anthropic

            client = Anthropic(api_key=self.api_key)

            presentation_text = "\n\n".join(
                [f"Slide {s['slide_number']}:\n{s['text']}" for s in self.slides_content]
            )

            message = client.messages.create(
                model="claude-opus-4-8",
                max_tokens=1024,
                messages=[
                    {
                        "role": "user",
                        "content": f"""Extract exactly {num_points} key points from this presentation.
Format each point as a single line, numbered 1-{num_points}.

Presentation:\n{presentation_text}"""
                    }
                ]
            )

            response_text = message.content[0].text
            points = [line.strip() for line in response_text.split('\n') if line.strip() and any(c.isdigit() for c in line[:3])]
            return points[:num_points]
        except Exception as e:
            logger.error(f"Error extracting key points: {e}")
            return []


class PPTXGenerator:
    """Generates new PPTX files from content."""

    def __init__(self, output_dir: str = "./output"):
        """
        Initialize PPTX generator.

        Args:
            output_dir: Directory for output files
        """
        self.output_dir = Path(output_dir)
        self.output_dir.mkdir(parents=True, exist_ok=True)

    def generate_from_text(self, content: str, title: str = "Generated Presentation") -> str:
        """Generate PPTX from text content."""
        try:
            from pptx.util import Inches, Pt
            from pptx.enum.text import PP_ALIGN

            prs = Presentation()
            prs.slide_width = Inches(10)
            prs.slide_height = Inches(7.5)

            title_slide_layout = prs.slide_layouts[0]
            slide = prs.slides.add_slide(title_slide_layout)
            title_shape = slide.shapes.title
            subtitle = slide.placeholders[1]

            title_shape.text = title
            subtitle.text = "Generated from content analysis"

            lines = content.split('\n')
            current_slide_lines = []

            for line in lines:
                if line.strip():
                    current_slide_lines.append(line)

                if len(current_slide_lines) >= 5 or (line.strip() == "" and current_slide_lines):
                    if current_slide_lines:
                        bullet_slide_layout = prs.slide_layouts[1]
                        slide = prs.slides.add_slide(bullet_slide_layout)
                        slide_title = slide.shapes.title
                        body_shape = slide.placeholders[1]

                        slide_title.text = "Key Points"

                        tf = body_shape.text_frame
                        tf.clear()

                        for line_text in current_slide_lines[:5]:
                            p = tf.add_paragraph()
                            p.text = line_text.strip()
                            p.level = 0

                        current_slide_lines = []

            output_file = self.output_dir / f"{title.replace(' ', '_')}.pptx"
            prs.save(str(output_file))
            logger.info(f"PPTX generated: {output_file}")
            return str(output_file)

        except Exception as e:
            logger.error(f"Error generating PPTX: {e}")
            return ""

    def generate_from_summary(self, summary: str, title: str = "Summary Presentation") -> str:
        """Generate PPTX from a summary text."""
        return self.generate_from_text(summary, title)


# ============================================================================
# MAIN WORKFLOW & CLI INTERFACE
# ============================================================================

def read_and_extract(pptx_path: str) -> Dict:
    """Read PPTX and extract content via API."""
    reader = PPTXReader(pptx_path)
    result = reader.read_pptx()

    if result.get("status") == "success" and reader.slides_content:
        preview = reader.slides_content[0].get("text", "")[:200] if reader.slides_content else ""
        return {
            "status": "success",
            "filename": result.get("filename"),
            "slides": result.get("slides"),
            "content_length": result.get("content_length"),
            "first_slide_preview": preview + "..." if len(preview) > 200 else preview,
        }
    elif result.get("status") != "success":
        return {
            "status": "error",
            "message": result.get("message", "Unknown error occurred"),
            "filename": None,
            "slides": None,
            "content_length": None
        }
    else:
        return {
            "status": "error",
            "message": "Failed to extract content from PPTX",
            "filename": result.get("filename"),
            "slides": result.get("slides")
        }


def analyze_and_answer(pptx_path: str, question: str) -> str:
    """Read PPTX and answer a question about it."""
    reader = PPTXReader(pptx_path)
    reader.read_pptx()

    if not reader.slides_content:
        return "Error: Could not extract content from PPTX"

    analyzer = DocumentAnalyzer(reader.slides_content)
    return analyzer.answer_question(question)


def generate_summary_pptx(pptx_path: str, output_path: Optional[str] = None) -> str:
    """Read PPTX, generate summary, and create new PPTX."""
    reader = PPTXReader(pptx_path)
    reader.read_pptx()

    if not reader.slides_content:
        return "Error: Could not extract content from PPTX"

    analyzer = DocumentAnalyzer(reader.slides_content)
    summary = analyzer.summarize(max_words=500)

    generator = PPTXGenerator(output_path or "./output")
    pptx_file = generator.generate_from_summary(
        summary,
        title=f"Summary of {Path(pptx_path).stem}"
    )

    if pptx_file:
        return f"Summary PPTX generated: {pptx_file}"
    return "Error: Failed to generate PPTX"


def main():
    """Main entry point for the skill."""
    import argparse

    parser = argparse.ArgumentParser(description="PPTX Content Extractor")
    subparsers = parser.add_subparsers(dest="command", help="Available commands")

    extract_parser = subparsers.add_parser("extract", help="Extract content from PPTX")
    extract_parser.add_argument("pptx_path", help="Path to PPTX file")

    answer_parser = subparsers.add_parser("answer", help="Answer question about PPTX")
    answer_parser.add_argument("pptx_path", help="Path to PPTX file")
    answer_parser.add_argument("question", help="Question to answer")

    summary_parser = subparsers.add_parser("summarize", help="Generate summary PPTX")
    summary_parser.add_argument("pptx_path", help="Path to PPTX file")
    summary_parser.add_argument("--output", help="Output directory")

    args = parser.parse_args()

    if args.command == "extract":
        result = read_and_extract(args.pptx_path)
        print(json.dumps(result, indent=2))
    elif args.command == "answer":
        answer = analyze_and_answer(args.pptx_path, args.question)
        print(answer)
    elif args.command == "summarize":
        result = generate_summary_pptx(args.pptx_path, args.output)
        print(result)
    else:
        parser.print_help()


if __name__ == "__main__":
    main()
